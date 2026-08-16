"""Render the W6 storyboard into a slide deck.

A title slide carrying the BLUF, then one slide per assertion beat: the headline
and its figure, nothing else. No bullets, no footers, no slide numbers — the
storyboard's own rule is that the sequence carries the argument with the slides
removed, so anything on the slide beyond the assertion and its evidence is
competing with the speaker.

Speaker notes are assembled from the storyboard's evidence column and from
``coursework/W6/talk_notes.md``; no note text is composed here.

The beats and their figures are transcribed from the storyboard table in
``coursework/W6/MSDS696_W6_Status_Report.md``, which is the authority; if a
headline is reworded there, reword it here. Figures come from ``img/`` and are
built by ``notebook/15_w6_visuals.ipynb`` — this script only lays them out, so
rebuilding a figure and re-running this is enough to update the deck.

Usage::

    python src/build_deck.py [-o OUTPUT]
"""

from __future__ import annotations

import argparse
from pathlib import Path
from typing import NamedTuple

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

REPO_ROOT = Path(__file__).resolve().parent.parent
IMG_DIR = REPO_ROOT / "img"
DEFAULT_OUTPUT = REPO_ROOT / "coursework" / "W6" / "MSDS696_W6_Deck.pptx"

# 16:9. The practice recording and the W8 final are both screen-shared.
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

MARGIN = Inches(0.6)
HEADLINE_TOP = Inches(0.45)
HEADLINE_H = Inches(1.35)
FIGURE_TOP = Inches(1.95)
FIGURE_BOTTOM_MARGIN = Inches(0.45)

INK = RGBColor(0x1A, 0x1A, 0x1A)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
HEADLINE_FONT = "Helvetica Neue"
HEADLINE_PT = 30
# Long assertions get a step down rather than a wrap onto a third line. Beat 12
# is the deck's longest headline at 108 characters.
HEADLINE_PT_LONG = 25
LONG_HEADLINE_CHARS = 72

# Title slide. It leads with the BLUF rather than the project title, which keeps
# the deck's assertion rule intact from the first slide — and makes the deck open
# and close on the same sentence, since beat 19 lands on it too. The project
# title, name and course sit beneath it as the identifying block.
TITLE_BLUF = "Stop targeting how big it gets.\nTarget where it starts."
TITLE_SUBTITLE = (
    "Predicting Region-Season Wildfire Cause Patterns "
    "to Target Prevention and Mitigation"
)
TITLE_BYLINE = "Craig Rudman  ·  MSDS 696 Practicum II  ·  Week 6"
TITLE_PT = 44
SUBTITLE_PT = 20
BYLINE_PT = 15
MUTED = RGBColor(0x5A, 0x5A, 0x5A)

class Beat(NamedTuple):
    """One storyboard beat and the speaker notes that go behind its slide.

    ``say`` is a line to deliver more or less as written and appears only where
    ``talk_notes.md`` already supplies one; ``evidence`` is the figures from the
    storyboard's evidence column; ``watch`` is a stated trap, caveat or framing
    obligation. **Nothing here is newly composed prose** — every line is lifted
    from the status report's storyboard table or from ``talk_notes.md``, so the
    notes cannot drift from the documents that are the authority for them.
    """

    headline: str
    figure: str
    say: str | None = None
    evidence: tuple[str, ...] = ()
    watch: tuple[str, ...] = ()


# Order is delivery order; see the storyboard table in the W6 status report.
BEATS: list[Beat] = [
    Beat(
        "Wildfires are seasonal.",
        "w6_seasonality.png",
        evidence=(
            "Most fires start in spring; most acres burn in summer.",
            "MAM and JJA start nearly the same number of fires and differ 3.9x in acres.",
        ),
        watch=("No axes and no magnitudes on this one — it is a calendar, not a chart.",),
    ),
    Beat(
        "Cause is regional, not national.",
        "w6_cause_map.png",
        evidence=(
            "All 105 Level III ecoregions, shaded by natural share of attributed acres.",
            "A West/East split at roughly the 100th meridian; Alaska almost entirely natural.",
            "Bimodal: 50 regions below 20% natural, 28 above 80%, only 27 in between.",
        ),
    ),
    Beat(
        "A region's cause mix is stable enough to forecast.",
        "w6_tier1_tiles.png",
        say="About three-quarters of the mix lands on the right cause.",
        evidence=(
            "Three tiles, worst to best: national average mix 42%, an even split 52%, "
            "the region's own history 73%.",
            "Acre-weighted TVD 0.580 / 0.485 / 0.266, forward-chained on 2010+, "
            "3,949 held-out region-seasons.",
        ),
        watch=(
            "Do NOT say \"right 73% of the time.\" 73% is an error magnitude (1 - TVD), "
            "not a hit rate — there is no \"of the time\" to attach it to.",
            "72.7% top-1 agreement is a different number that DOES mean \"of the time.\" "
            "Only use that phrasing if the slide on screen shows top-1.",
            "The middle tile carries the weight: the national mix at 42% is worse than "
            "guessing, which is beat 2 reappearing as forecast error.",
        ),
    ),
    Beat(
        "A few years of rolling average is enough.",
        "w6_k_sweep.png",
        evidence=(
            "Acre-weighted TVD falls 0.331 -> 0.278 from one prior season to three, then flattens.",
            "Every window from three up sits within 1.4 points of the best.",
            "Cause composition is a standing property, not a yearly swing.",
        ),
        watch=(
            "The compression candidate if the talk runs long — it tunes a parameter of "
            "beat 3's baseline rather than adding a claim.",
        ),
    ),
    Beat(
        "Within Human, history names the lead cause more often than not.",
        "w6_human_tiles.png",
        evidence=(
            "Same three tiles as beat 3, one level deeper: how often the predicted "
            "leading human sub-cause is right out of 11.",
            "An even split 9%, the national human mix 16%, the region's own history 54%.",
            "Acre-weighted TVD 0.489 against the national mix's 0.643; 3,850 held-out "
            "region-seasons.",
        ),
        watch=(
            "54% is the Human floor on Human's own population. Do NOT quote it as an "
            "end-to-end number — end to end is 46.2%.",
        ),
    ),
    Beat(
        "A learned model made it worse.",
        "w6_human_ladder.png",
        evidence=(
            "Three rungs against the floor line: the region's own history names the "
            "leading human cause 54% of the time; gradient boosting on region character 36%; "
            "the same model given that history as a feature 47%.",
            "Acre-weighted TVD 0.489 / 0.588 / 0.554.",
            "Even handed the winning quantity, the model cannot beat taking its mean.",
        ),
        watch=(
            "Concede the fair part: the rungs were run once at standard settings and not "
            "tuned. What the result rules out is \"you never gave it the right features.\"",
            "Keep this beat — it is the only one conceding a model was tried and lost, "
            "which is what keeps the nulls credible.",
        ),
    ),
    Beat(
        "Where fires start is predictable, but not at ecoregion scale.",
        "w6_siting_glance.png",
        say=(
            "Everything so far has been a whole ecoregion — one number for an area the "
            "size of a small state. That is the right scale for deciding what to target, "
            "and the wrong one for deciding where to put anything. So from here the map "
            "breaks into cells of about sixty thousand acres, and the question changes "
            "with it: not how much will burn, but where fires start."
        ),
        evidence=(
            "Klamath hexes in two bands: the deep band is 6% of the region catching 32% "
            "of next season's starts (5.2x); the light band 29% for 60% (2.1x).",
            "Held-out Spearman: Human +0.53, Natural +0.34.",
        ),
        watch=(
            "The deck's ONE grain change, and three things change at once: unit "
            "(105 ecoregions -> 36,234 hexes), target (acres -> ignition counts), "
            "answer (shares -> counts). The headline announces only the first; say the others.",
            "This is NOT an acres model. Unsaid, the audience reads the capture curve as "
            "\"32% of the burn under 6% of the ground\" — a much stronger claim than the "
            "one being made. It is 32% of the STARTS.",
            "The return decays fast: 90% of starts needs 77.8% of the ground at 1.16x. "
            "The ranking concentrates return; it does not eliminate the tail.",
            "Q&A companion: img/w6_capture_curve.png.",
        ),
    ),
    Beat(
        "That skill is spatial, not statistical luck.",
        "w6_shuffled_control.png",
        say=(
            "I am not changing how many fires I predict. I am only changing where I say "
            "they will be. Same numbers, wrong places — and the prediction stops working."
        ),
        evidence=(
            "Two lines, forecast and shuffled: mean observed starts against mean predicted "
            "for 20 equal-count strata.",
            "The forecast climbs the diagonal (top stratum predicted 4.7, observed 3.7); "
            "the same predictions dealt to the wrong hexes go flat at ~0.4.",
            "Spearman +0.526 -> +0.0002; MAE 0.43 -> 0.77, worse than the uniform "
            "baseline's 0.70. 1.59M held-out hex-seasons.",
        ),
        watch=(
            "Avoid \"the model is accurate\" — the line runs below the diagonal and "
            "under-predicts the busiest hexes. Say: it ranks well, it does not promise counts.",
            "The shuffled line is FLAT, not low. It sits at ~0.4, the average across all hexes.",
            "Why shuffled and not random: shuffling changes exactly one thing (the pairing), "
            "so the collapse is attributable to siting alone.",
        ),
    ),
    Beat(
        "Human fire is predictable year-round; lightning only in summer.",
        "w6_season_skill.png",
        evidence=(
            "Held-out Spearman by season, each branch scored separately in all 11 held-out "
            "years, band spanning the observed year-to-year range.",
            "Human runs flat and high (median +0.47 to +0.61, peaking in spring); natural "
            "is a summer surface (+0.42 JJA, +0.07 DJF).",
            "Human beats natural in all 44 season-years without exception.",
        ),
        watch=("The deck's only figure showing a distribution rather than a point estimate.",),
    ),
    Beat(
        "Almost all the acres are in almost none of the cells.",
        "w6_acres_concentration.png",
        evidence=(
            "Cumulative share of natural acres, cells ordered by acres burned, least to most.",
            "The worst-burning 10% of cells hold 98% of the acres; the worst-burning 1% hold 55%.",
        ),
        watch=(
            "This sets the stakes for beat 11 — it names which cells a forecast has to get "
            "right BEFORE any error is shown. Matched pair with 11; neither survives alone.",
        ),
    ),
    Beat(
        "Up to a point, both are predictable. Past that point, natural fire is harder.",
        "w6_branch_deciles.png",
        evidence=(
            "Typical forecast error against how much a cell burned, both branches on one axis — "
            "the same axis and direction as beat 10, so the right edge means \"the big burns\" "
            "on both slides.",
            "They track through the smallest third, then separate: natural runs 2-3x worse than "
            "human at the same cell size, reaching 687x on a median 8,061-acre cell against "
            "human's 19x on 240.",
            "Cells under 1 acre excluded: 25.3% of FPA-FOD rows sit at exactly 0.1 acres "
            "(44.5% of natural fires), a reporting default rather than a measurement.",
        ),
        watch=(
            "Quote the population with the number. Across ALL JJA natural burning cells the "
            "top decile is 269.8x on a median 2,970-acre cell; 854.9x is the six-forest-"
            "ecoregion population used for the covariate ladder. The like-for-like against "
            "Human is 269.8x, not 855x.",
            "This licenses shipping two DIFFERENT products: Human can be ranked by expected "
            "acres, Natural cannot.",
        ),
    ),
    Beat(
        "We tried to fix it with drought and fuel, but where fires start is a "
        "property of the place, not of the year.",
        "w6_ignition_ladder.png",
        say=(
            "Raw: do greener hexes have more fires than browner hexes? Yes, moderately. "
            "Within-hex: when a hex is greener than its own normal, does it have more fires "
            "than its own normal? Barely. The first question is answered by which hex you are "
            "looking at; the second by which year it is — and a forecast needs the second."
        ),
        evidence=(
            "Two flat lines, one per branch, across the rungs: the region's own history, "
            "+ drought, + fuel load, + both. Nothing moves.",
            "Best gain on either branch is +0.0045; the y-axis runs from zero so a real "
            "effect would have been visible.",
            "Both branches on the same 29,293 held-out cells (JJA, the six forest ecoregions).",
            "The measured reason, spoken not drawn: pdsi -0.137 -> -0.073 and "
            "NDVI +0.228 -> +0.098 from raw to within-hex anomaly.",
        ),
        watch=(
            "Frame as a REPAIR ATTEMPT, not a new topic — beat 11 leaves a failure on the "
            "table and beats 12-14 are what was done about it.",
            "Resist \"fires happen where the fuel is\" — fuel state added +0.004. The correct "
            "compression is \"fires happen where fires have happened.\"",
            "Q&A companion: img/w6_ndvi_variance.png, the place-vs-year split at 2.8x.",
        ),
    ),
    Beat(
        "The same data does predict how much burns.",
        "w6_acres_ladder.png",
        evidence=(
            "Deliberately the same figure as beat 12: same rungs, same axis, same zero-based "
            "scale, one line instead of two. The shape is the argument.",
            "Climate + NDVI together +0.049, 26.6 SD above a covariate-shuffled control, "
            "holding across five forward-chaining split years (+0.012 to +0.066).",
            "Neither half works alone — drought alone -0.008, fuel alone +0.001. Wet heavy "
            "fuel will not carry fire; dry bare ground has nothing to burn.",
        ),
        watch=(
            "Different target and population from beat 12 (7,799 burning JJA cells, "
            "burn-conditional baseline) — comparable in SHAPE, not cell for cell.",
        ),
    ),
    Beat(
        "But the gain lands on the fires nobody needed predicted.",
        "w6_gain_landing.png",
        evidence=(
            "Beat 11's axis again, with the covariate model laid over the baseline.",
            "The shaded lens opens in deciles 6-9 (1-200 acre cells) and closes at the right "
            "edge where the two lines finish together: 855x -> 868x on a median 5,073-acre cell.",
            "Deciles 1-5 get materially worse — decile 1 nearly doubles, 18.8x -> 35.9x. "
            "780 held-out cells per decile.",
        ),
        watch=(
            "Beat 10 is what makes this fatal rather than disappointing: the improvement "
            "misses the cells holding 98% of the acres.",
            "Third beat landing the eye on the same right edge — the argument for reframing "
            "rather than tuning.",
        ),
    ),
    Beat(
        "Siting needed a finer place. Size needs a finer moment.",
        "w6_grain_parallel.png",
        say=(
            "Remember what we did to make siting work: the region was too coarse to put "
            "anything anywhere, so we dropped down to a hex. This is the same problem on the "
            "other axis. We asked how big a fire gets over a whole season, and a season is too "
            "coarse a unit to answer that — what makes a fire run is the wind on a particular "
            "afternoon, whether crews were already committed, what time of day it started. We "
            "could not test that here, because same-day data is a different project. But the "
            "shape of the failure tells you where to look."
        ),
        evidence=(
            "The deck's only figure that plots no data, and its only one about the method "
            "rather than the fire. The solved row is closed, the untested row dashed.",
            "The fourth time the project has met the same lesson: W4's pooled climate null, "
            "beat 12's places-not-years, beat 7's grain drop, now this.",
        ),
        watch=(
            "Say \"before the season\" whenever the claim is stated as a null. \"Megafire size "
            "is unpredictable\" is a much bigger claim than anything measured here.",
            "Do not apologise for the null. Five ablations, a shuffled control at 26.6 SD, and "
            "a gain that landed in the wrong deciles is a thorough negative result.",
            "If asked \"did you try hard enough?\" — pre-season data was tried hard, same-day "
            "data was not tried at all. The second half is the open question, not a gap.",
            "If the room needs \"stop targeting it\" said aloud, say it as the last sentence "
            "rather than the headline.",
        ),
    ),
    Beat(
        "Every megafire was an ignition first.",
        "w6_ignition_gate.png",
        evidence=(
            "Two bars in a common frame, no axis: a hex-season that ignites at all produces a "
            ">=1,000-acre burn 6.7% of the time against 0.29% for one that does not — a 22.8x gate.",
            "JJA natural, held-out years. Percentages printed because at true scale the 0.29% "
            "bar is nearly invisible — which is itself the finding.",
        ),
        watch=(
            "Do not oversell the gate: it is necessary, not sufficient. 93% of igniting "
            "hex-seasons still produce nothing large. It narrows the field; it does not "
            "identify the fire.",
            "This is the door beat 15 left open — size is not forecastable, but the event "
            "upstream of it is.",
        ),
    ),
    Beat(
        "One ignition is enough.",
        "w6_one_is_enough.png",
        evidence=(
            "One stacked bar over all 2,724 large-fire cells in the held-out years, split by "
            "how many times their hex ignited that season: 49% had exactly one, 21% two, "
            "30% three or more.",
            "Only the first segment carries color — it is the ground a planner would "
            "deprioritise by ranking on ignition count.",
        ),
        watch=(
            "The rule is binary: does this place ignite, not how often.",
            "If asked \"doesn't a cell with more ignitions carry more risk?\" — the rate does "
            "rise (19.1% at 11-20 ignitions vs 5.4% at one), but risk PER ignition falls "
            "0.054 -> 0.014, and 49% of large-fire cells had exactly one ignition.",
            "Ignition count ranks burned area worse than the hex's own burn history does: "
            "Spearman +0.253 against +0.357.",
        ),
    ),
    Beat(
        "Nearly a fifth of burned acres have no specific cause — and that gap is "
        "itself forecastable.",
        "w6_unknown_triage.png",
        evidence=(
            "The ranked triage list, headed by Southwestern Tablelands MAM at 1.17M predicted "
            "unattributed acres.",
            "Ranked by acres rather than by rate: Central Great Plains has the worse "
            "attribution rate (66%) but a fifth of the burn.",
            "Unknown-branch persistence: acre-weighted MAE 0.167 against the global mean's 0.240.",
        ),
        watch=(
            "The third leg of the recommendation, delivered where it is actionable rather "
            "than as a finding in its own right.",
        ),
    ),
    Beat(
        "Target causes by region. Site the work by ignition. Fix the record where "
        "it says neither.",
        "w6_recommendation.png",
        say=("Stop targeting how big it gets. Target where it starts."),
        evidence=(
            "Three rows, one per Tier-1 class, ordered by share so the top row is where the "
            "acres are: Natural 58.9%, site by ignition, hex-season; Human 22.7%, rank causes "
            "by the acres they drive, ecoregion-season; Unknown 18.5%, fix the record "
            "underneath, ecoregion-season.",
            "The boundary, on the line beneath: rank on all three — the order is trustworthy "
            "and the acre level much less so.",
        ),
        watch=(
            "A TARGETING claim, not an efficacy one. Nothing here measures what a treatment "
            "achieves — no before/after, no control, no counterfactual.",
            "If asked \"so if we treat those hexes, we cut the burn?\" — say no, plainly. "
            "Davis et al. (2024) on severity reduction is separate evidence from someone "
            "else's study; cite it if asked, do not fold it in.",
            "If they remember one thing: site the work against where fires start, because "
            "that is the one stage of the escalation this data can see in advance.",
        ),
    ),
]

# Beats in the five-minute practice cut. Kept as data rather than a second beat
# list so the cut stays a subset of one storyboard, which is how the report
# argues it: 1, 2, 3, 7, 10, 15, 16, 17, 19.
FIVE_MINUTE_CUT = (1, 2, 3, 7, 10, 15, 16, 17, 19)


def add_headline(slide, text: str) -> None:
    """Place the assertion across the top of the slide, left-aligned."""
    box = slide.shapes.add_textbox(
        MARGIN, HEADLINE_TOP, SLIDE_W - 2 * MARGIN, HEADLINE_H
    )
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP

    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(HEADLINE_PT_LONG if len(text) > LONG_HEADLINE_CHARS else HEADLINE_PT)
    run.font.bold = True
    run.font.color.rgb = INK
    run.font.name = HEADLINE_FONT


def add_figure(slide, image_path: Path) -> None:
    """Fit the figure inside the region below the headline, centered.

    Scaled to fit rather than to fill: the figures range from 1.42:1 to 2.85:1,
    so a fixed frame would crop the wide ones and every one of them carries
    printed labels at its edges.
    """
    frame_w = SLIDE_W - 2 * MARGIN
    frame_h = SLIDE_H - FIGURE_TOP - FIGURE_BOTTOM_MARGIN

    with Image.open(image_path) as im:
        px_w, px_h = im.size

    scale = min(frame_w / px_w, frame_h / px_h)
    width = Emu(int(px_w * scale))
    height = Emu(int(px_h * scale))
    left = Emu(int(MARGIN + (frame_w - width) / 2))
    top = Emu(int(FIGURE_TOP + (frame_h - height) / 2))

    slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)


def add_title_slide(prs, layout) -> None:
    """Open the deck on the BLUF, with the identifying block beneath it.

    Vertically centered as one stack rather than pinned to the top, so the slide
    reads as a statement rather than as a header with metadata under it.
    """
    slide = prs.slides.add_slide(layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER

    box = slide.shapes.add_textbox(
        MARGIN, Inches(2.1), SLIDE_W - 2 * MARGIN, Inches(3.3)
    )
    frame = box.text_frame
    frame.word_wrap = True

    # (text, point size, bold, color, space-before in points)
    lines = [
        (TITLE_BLUF.split("\n")[0], TITLE_PT, True, INK, 0),
        (TITLE_BLUF.split("\n")[1], TITLE_PT, True, INK, 0),
        (TITLE_SUBTITLE, SUBTITLE_PT, False, MUTED, 30),
        (TITLE_BYLINE, BYLINE_PT, False, MUTED, 18),
    ]
    for i, (text, size, bold, color, space_before) in enumerate(lines):
        paragraph = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        paragraph.alignment = PP_ALIGN.LEFT
        if space_before:
            paragraph.space_before = Pt(space_before)
        run = paragraph.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = HEADLINE_FONT

    slide.notes_slide.notes_text_frame.text = (
        "SAY —\n"
        "Stop targeting how big it gets. Target where it starts.\n\n"
        "WATCH —\n"
        "- A TARGETING claim, not an efficacy one. Nothing in this project measures "
        "what a treatment achieves; the ranking says where fire is most likely to "
        "arrive, which is a necessary condition for sited work to pay off and not a "
        "sufficient one.\n"
        "- The deck closes on this same sentence at beat 19, so first and last "
        "statements match."
    )


def add_notes(slide, beat: Beat) -> None:
    """Write the beat's speaker notes into the slide's notes page.

    Laid out as SAY / EVIDENCE / WATCH rather than as a paragraph, because the
    notes pane is read in glances while presenting. SAY comes first and only
    exists where ``talk_notes.md`` supplies a line; the other two are lists so
    the eye can find one number without reading the block.
    """
    if not (beat.say or beat.evidence or beat.watch):
        return

    frame = slide.notes_slide.notes_text_frame
    blocks: list[str] = []
    if beat.say:
        blocks.append(f"SAY —\n{beat.say}")
    if beat.evidence:
        blocks.append("EVIDENCE —\n" + "\n".join(f"- {line}" for line in beat.evidence))
    if beat.watch:
        blocks.append("WATCH —\n" + "\n".join(f"- {line}" for line in beat.watch))

    frame.text = "\n\n".join(blocks)


def build(output_path: Path, beats: list[Beat]) -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]

    add_title_slide(prs, blank_layout)

    for beat in beats:
        image_path = IMG_DIR / beat.figure
        if not image_path.exists():
            raise FileNotFoundError(f"storyboard figure missing: {image_path}")

        slide = prs.slides.add_slide(blank_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = PAPER

        add_headline(slide, beat.headline)
        add_figure(slide, image_path)
        add_notes(slide, beat)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("-o", "--output", type=Path, default=DEFAULT_OUTPUT)
    parser.add_argument(
        "--five-minute",
        action="store_true",
        help="build only the nine beats of the practice-recording cut",
    )
    args = parser.parse_args()

    beats = BEATS
    output = args.output
    if args.five_minute:
        beats = [BEATS[n - 1] for n in FIVE_MINUTE_CUT]
        if output == DEFAULT_OUTPUT:
            output = output.with_name("MSDS696_W6_Deck_5min.pptx")

    written = build(output, beats)
    print(f"title + {len(beats)} beats -> {written.relative_to(REPO_ROOT)}")


if __name__ == "__main__":
    main()
